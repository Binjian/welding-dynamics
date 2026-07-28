"""Build a PowerPoint deck from MJ Warp, MJX, Physx.md.

The deck reuses the theme/master from NVIDIA_四大家族_合作调研报告.pptx and
recreates its visual language with native, editable PowerPoint shapes.

Run with:
    PYTHONPATH=/usr/lib/python3/dist-packages python docs/build_mjwarp_mjx_physx_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "NVIDIA_四大家族_合作调研报告.pptx"
OUTPUT = ROOT / "MJWarp_MJX_PhysX_机器人仿真选型.pptx"

FONT = "Noto Sans CJK SC"
FONT_MONO = "Noto Sans Mono CJK SC"

# Palette sampled from the supplied template.
INK = "262626"
MUTED = "727171"
BLUE = "026BFF"
NAVY = "0A2350"
CYAN = "19B7FF"
GREEN = "24C58D"
ORANGE = "F68941"
GOLD = "E7BE5D"
PURPLE = "926DE5"
WHITE = "FFFFFF"
PALE = "F4F7FC"
PALE_BLUE = "EAF2FF"
PALE_GREEN = "E7F8F1"
PALE_ORANGE = "FDEEE0"
PALE_PURPLE = "F2EDFC"
BORDER = "E4E9F2"
NAVY_LINE = "2A417A"
LIGHT_TEXT = "C4CFE6"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str | None):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str | None, width: float = 1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def rect(
    slide,
    x,
    y,
    w,
    h,
    fill=WHITE,
    line=None,
    radius=True,
    line_width=1.0,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, line_width)
    # Reduce the default corner radius for wide cards.
    if radius and shape.adjustments:
        shape.adjustments[0] = 0.08
    return shape


def circle(slide, x, y, d, fill, line=None, line_width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    set_fill(shape, fill)
    set_line(shape, line, line_width)
    return shape


def line(slide, x1, y1, x2, y2, color=BORDER, width=1.0):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.0,
    fit=False,
    rotation=0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.rotation = rotation
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    if fit:
        tf.fit_text(font_family=font, max_size=Pt(size))
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = str(value)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def rich_text(
    slide,
    runs,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    for item in runs:
        r = p.add_run()
        r.text = item["text"]
        r.font.name = item.get("font", FONT)
        r.font.size = Pt(item.get("size", size))
        r.font.bold = item.get("bold", False)
        r.font.color.rgb = rgb(item.get("color", color))
    return box


def label_pill(slide, label, x, y, w, color, *, size=10.5):
    rect(slide, x, y, w, 0.36, WHITE, color, radius=True, line_width=1)
    circle(slide, x + 0.13, y + 0.105, 0.15, color)
    text(slide, label, x + 0.36, y, w - 0.43, 0.36, size=size, color=INK)


def section_header(slide, section, eyebrow, title_value, *, color=BLUE):
    rect(slide, 0.9, 0.62, 0.62, 0.62, color, None, radius=False)
    text(
        slide,
        section,
        0.9,
        0.62,
        0.62,
        0.62,
        size=19,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    text(slide, eyebrow, 1.68, 0.6, 10.5, 0.32, size=10.5, color=color, bold=True)
    text(slide, title_value, 1.66, 0.91, 10.75, 0.58, size=27, color=INK, bold=True)


def footer(slide, page):
    text(
        slide,
        "MJLAB / PHYSICS BACKEND REVIEW",
        0.9,
        7.12,
        4.0,
        0.18,
        size=7.5,
        color="9AA4B5",
        valign=MSO_ANCHOR.BOTTOM,
    )
    text(
        slide,
        f"{page:02d}",
        11.9,
        7.08,
        0.5,
        0.22,
        size=8,
        color="9AA4B5",
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.BOTTOM,
    )


def bullet_lines(
    slide,
    items,
    x,
    y,
    w,
    row_h,
    *,
    dot_color=BLUE,
    size=12,
    text_color=MUTED,
    bold_first=False,
):
    for i, item in enumerate(items):
        yy = y + i * row_h
        circle(slide, x, yy + row_h * 0.42 - 0.055, 0.11, dot_color)
        text(
            slide,
            item,
            x + 0.25,
            yy,
            w - 0.25,
            row_h,
            size=size,
            color=text_color,
            bold=(bold_first and i == 0),
            valign=MSO_ANCHOR.MIDDLE,
        )


def stage_arrow(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    set_line(shape, None)
    if shape.adjustments:
        shape.adjustments[0] = 0.22
    return shape


def wipe_template_slides(prs: Presentation):
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    circle(slide, 9.9, -2.5, 6.8, PALE_BLUE)
    circle(slide, 11.2, 3.7, 3.5, PALE, CYAN, 0.7)
    # Small abstract physics/network mark.
    circle(slide, 11.62, 0.74, 1.0, BLUE)
    for xx, yy in [(11.91, 0.98), (12.23, 0.98), (12.07, 1.29)]:
        circle(slide, xx, yy, 0.14, WHITE)
    line(slide, 11.98, 1.05, 12.26, 1.05, WHITE, 1.3)
    line(slide, 12.00, 1.10, 12.12, 1.31, WHITE, 1.3)
    line(slide, 12.28, 1.10, 12.15, 1.31, WHITE, 1.3)

    rect(slide, 0.9, 1.62, 0.55, 0.09, BLUE, None, radius=False)
    text(
        slide,
        "NVIDIA PHYSICS STACK · ROBOTICS SIMULATION",
        0.9,
        1.78,
        9.8,
        0.38,
        size=12.5,
        color=BLUE,
        bold=True,
    )
    text(slide, "MJWarp · MJX · PhysX", 0.86, 2.25, 11.5, 0.84, size=42, bold=True)
    text(
        slide,
        "机器人仿真后端对比、平台应用与 MjLab 选型建议",
        0.88,
        3.2,
        11.2,
        0.72,
        size=27,
        color=BLUE,
        bold=True,
    )
    text(
        slide,
        "从计算架构、接触保真度到训练—验证工作流",
        0.9,
        4.12,
        10.5,
        0.42,
        size=14,
        color=MUTED,
    )
    label_pill(slide, "MJX · JAX/XLA", 0.9, 4.82, 1.68, PURPLE)
    label_pill(slide, "MJWarp · CUDA", 2.82, 4.82, 1.82, GREEN)
    label_pill(slide, "PhysX · Isaac", 4.88, 4.82, 1.7, ORANGE)
    label_pill(slide, "MjLab", 6.82, 4.82, 1.05, CYAN)
    text(
        slide,
        "基于《MJ Warp, MJX, Physx.md》与指定模板整理",
        0.9,
        6.63,
        11,
        0.3,
        size=9.5,
        color=MUTED,
    )
    return slide


def add_contents_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    text(slide, "目录  CONTENTS", 0.9, 0.7, 8, 0.7, size=28, bold=True)
    rect(slide, 0.92, 1.42, 0.6, 0.08, BLUE, None, radius=False)
    rows = [
        ("01", "核心结论与技术架构", "三种后端不是互斥替代，而是面向不同任务的组合", BLUE),
        ("02", "接触物理与性能差异", "动态分支、接触求解、摩擦模型与几何复杂度", GREEN),
        ("03", "人形机器人训练工作流", "MJX 训练 → MJWarp 验证 → 真实机器人", GOLD),
        ("04", "PhysX 对比与引擎选型", "任务形态 · 硬件平台 · 工程生态", ORANGE),
        ("05", "MjLab 路线与应用附录", "能力建设 · 新型机器人 · 焊接机器人 · 多物理场仿真", CYAN),
    ]
    for i, (num, title_value, sub, color) in enumerate(rows):
        yy = 1.82 + i * 0.98
        rect(slide, 0.9, yy, 11.5, 0.86, PALE, None, radius=False)
        text(slide, num, 1.15, yy, 1.2, 0.86, size=31, color=color, bold=True)
        rect(slide, 2.55, yy + 0.18, 0.03, 0.5, color, None, radius=False)
        text(slide, title_value, 2.8, yy + 0.06, 8.9, 0.37, size=15.5, bold=True)
        text(slide, sub, 2.8, yy + 0.42, 9.1, 0.28, size=10.5, color=MUTED)
    footer(slide, 2)
    return slide


def add_executive_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "01", "DECISION  / 核心结论", "最佳答案不是单选，而是一条分层工作流")

    rect(slide, 0.9, 1.85, 3.45, 4.85, NAVY, None)
    text(slide, "1", 0.9, 2.24, 3.45, 1.0, size=68, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "套组合策略", 0.9, 3.15, 3.45, 0.42, size=19, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    text(
        slide,
        "训练吞吐量\n×\n物理保真度\n×\n工程生态",
        1.18,
        3.82,
        2.9,
        1.7,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    line(slide, 1.5, 5.75, 3.75, 5.75, NAVY_LINE, 1)
    text(
        slide,
        "按任务阶段切换后端，\n比押注单一引擎更稳妥",
        1.1,
        5.91,
        3.05,
        0.6,
        size=11.5,
        color=LIGHT_TEXT,
        align=PP_ALIGN.CENTER,
    )

    cards = [
        ("MJX", "训练引擎", "大规模并行 · JAX 生态\n运动策略 / 简单接触", PURPLE, PALE_PURPLE, "01"),
        ("MJWarp", "验证与复杂交互", "动态接触 · 复杂几何\n灵巧操作 / 高保真验证", GREEN, PALE_GREEN, "02"),
        ("PhysX", "工程与大场景", "Isaac / Omniverse 生态\n传感器 · 资产 · 通用仿真", ORANGE, PALE_ORANGE, "03"),
    ]
    for i, (name, role, desc, color, pale, num) in enumerate(cards):
        yy = 1.85 + i * 1.55
        rect(slide, 4.7, yy, 7.7, 1.37, WHITE, BORDER, radius=True)
        rect(slide, 4.98, yy + 0.35, 0.64, 0.64, color, None, radius=False)
        text(slide, num, 4.98, yy + 0.35, 0.64, 0.64, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, name, 5.88, yy + 0.19, 2.05, 0.42, size=19, color=INK, bold=True)
        rect(slide, 7.95, yy + 0.21, 1.72, 0.34, pale, None)
        text(slide, role, 8.04, yy + 0.21, 1.54, 0.34, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, 5.9, yy + 0.64, 6.05, 0.56, size=11.5, color=MUTED)
    rect(slide, 4.7, 6.5, 7.7, 0.2, BLUE, None, radius=False)
    footer(slide, 3)
    return slide


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "02", "ARCHITECTURE  / 底层计算", "计算范式决定了各自擅长的物理问题")

    cols = [
        (0.9, "MJX", "JAX / XLA", PURPLE, PALE_PURPLE, "静态计算图", [
            "大批量同构环境",
            "固定大小数据结构",
            "TPU / GPU / Mac 可用",
        ], "优势：向量化吞吐与 JAX 工作流"),
        (4.82, "MJWarp", "NVIDIA Warp / CUDA", GREEN, PALE_GREEN, "SIMT 动态内核", [
            "线程可独立分支",
            "仅计算活跃接触",
            "面向 NVIDIA GPU",
        ], "优势：接触丰富场景与复杂几何"),
        (8.74, "PhysX 5", "Isaac Sim / Lab", ORANGE, PALE_ORANGE, "通用实时引擎", [
            "TGS 约束求解",
            "大型场景与传感器",
            "Omniverse 资产生态",
        ], "优势：成熟工具链与工程集成"),
    ]
    for x, name, stack, color, pale, core, bullets, verdict in cols:
        rect(slide, x, 1.86, 3.66, 4.75, WHITE, color, radius=True, line_width=1.4)
        circle(slide, x + 0.25, 2.12, 0.62, color)
        text(slide, name[0], x + 0.25, 2.12, 0.62, 0.62, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, name, x + 1.02, 2.02, 2.35, 0.42, size=20, bold=True)
        text(slide, stack, x + 1.02, 2.43, 2.35, 0.3, size=10.5, color=color, bold=True)
        rect(slide, x + 0.25, 3.02, 3.16, 0.72, pale, None)
        text(slide, core, x + 0.25, 3.02, 3.16, 0.72, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
        bullet_lines(slide, bullets, x + 0.35, 3.95, 2.95, 0.56, dot_color=color, size=11.2)
        line(slide, x + 0.25, 5.78, x + 3.41, 5.78, BORDER, 1)
        text(slide, verdict, x + 0.28, 5.9, 3.1, 0.5, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 4)
    return slide


def add_mjx_mjwarp_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "02", "PERFORMANCE  / MJX vs MJWarp", "接触越复杂，动态计算的价值越明显")

    rect(slide, 0.9, 1.82, 3.15, 2.05, NAVY, None)
    text(slide, "50–100×", 1.05, 2.05, 2.85, 0.75, size=43, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "复杂操作任务", 1.05, 2.78, 2.85, 0.35, size=15, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "原文给出的 MJWarp 相对 MJX\n性能量级（取决于场景）", 1.15, 3.18, 2.65, 0.48, size=9.8, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

    rect(slide, 0.9, 4.05, 3.15, 2.45, PALE, None)
    text(slide, "4000+", 1.05, 4.35, 2.85, 0.65, size=38, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "并行人形环境", 1.05, 5.02, 2.85, 0.35, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "原文用于说明 MJX 在\n运动策略训练中的吞吐优势", 1.15, 5.52, 2.65, 0.55, size=10.3, color=MUTED, align=PP_ALIGN.CENTER)

    x0, y0 = 4.35, 1.82
    widths = [2.05, 2.95, 2.95]
    headers = ["维度", "MJX", "MJWarp"]
    colors = [BLUE, PURPLE, GREEN]
    xx = x0
    for w, hdr, color in zip(widths, headers, colors):
        rect(slide, xx, y0, w, 0.55, color, None, radius=False)
        text(slide, hdr, xx, y0, w, 0.55, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        xx += w
    rows = [
        ("接触规模", "少且可预测", "多且动态"),
        ("数据结构", "固定 / 填充", "按活跃接触计算"),
        ("复杂网格", "扩展性较弱", "更适合复杂几何"),
        ("主要任务", "移动 / 步态训练", "操作 / 装配 / 抓取"),
        ("关键平台", "JAX · TPU · GPU", "CUDA · NVIDIA GPU"),
    ]
    for i, row in enumerate(rows):
        yy = y0 + 0.66 + i * 0.78
        xx = x0
        fills = [PALE, WHITE, WHITE]
        for j, (w, val) in enumerate(zip(widths, row)):
            rect(slide, xx, yy, w, 0.68, fills[j], BORDER, radius=False, line_width=0.7)
            text(
                slide,
                val,
                xx + 0.08,
                yy,
                w - 0.16,
                0.68,
                size=11.5 if j else 11,
                color=INK if j == 0 else MUTED,
                bold=(j == 0),
                align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,
            )
            xx += w
    rect(slide, 4.35, 6.44, 8.0, 0.24, BLUE, None, radius=False)
    footer(slide, 5)
    return slide


def draw_contact_panel(slide, x, y, w, title_value, subtitle, color, hard):
    rect(slide, x, y, w, 4.05, WHITE, color, radius=True, line_width=1.3)
    circle(slide, x + 0.28, y + 0.28, 0.62, color)
    text(slide, "H" if hard else "S", x + 0.28, y + 0.28, 0.62, 0.62, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, title_value, x + 1.04, y + 0.22, w - 1.3, 0.42, size=19, bold=True)
    text(slide, subtitle, x + 1.04, y + 0.66, w - 1.3, 0.28, size=10.5, color=color, bold=True)

    # Stylized foot and ground.
    foot_y = y + 1.55 if hard else y + 1.7
    rect(slide, x + 1.22, foot_y, w - 2.44, 0.62, PALE, color, radius=True, line_width=1.3)
    text(slide, "机器人脚 / 刚体", x + 1.22, foot_y, w - 2.44, 0.62, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    if hard:
        rect(slide, x + 0.66, y + 2.34, w - 1.32, 0.18, NAVY, None, radius=False)
        # Impulse marker.
        impulse = slide.shapes.add_shape(
            MSO_SHAPE.UP_ARROW,
            Inches(x + w / 2 - 0.22),
            Inches(y + 2.03),
            Inches(0.44),
            Inches(0.56),
        )
        set_fill(impulse, color)
        set_line(impulse, None)
        text(slide, "d = 0  →  瞬时冲量", x + 0.65, y + 2.62, w - 1.3, 0.4, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        detail = "严格非穿透 · 约束求解\n冲击与振动更接近真实硬地面"
    else:
        # Spring coils as short alternating lines.
        for k in range(5):
            sx = x + w / 2 - 0.38 + k * 0.16
            line(slide, sx, y + 2.28 + (0.08 if k % 2 else 0), sx + 0.16, y + 2.36 - (0.08 if k % 2 else 0), color, 1.8)
        rect(slide, x + 0.66, y + 2.48, w - 1.32, 0.18, NAVY, None, radius=False)
        text(slide, "d < 0  →  F ≈ k·d", x + 0.65, y + 2.77, w - 1.3, 0.4, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        detail = "允许微小穿透 · 罚函数\n地面像微观弹簧，训练更平滑"
    text(slide, detail, x + 0.55, y + 3.18, w - 1.1, 0.62, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_contact_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "02", "CONTACT PHYSICS  / 接触保真度", "“蹦床效应”来自软接触，“硬着陆”来自约束求解")
    draw_contact_panel(slide, 0.9, 1.82, 5.53, "MJX · 软接触", "平滑、静态、训练友好", PURPLE, hard=False)
    draw_contact_panel(slide, 6.87, 1.82, 5.53, "MJWarp · 硬接触", "动态、严格、验证友好", GREEN, hard=True)

    rect(slide, 0.9, 6.08, 11.5, 0.62, PALE_BLUE, None)
    text(slide, "摩擦模型", 1.15, 6.08, 1.25, 0.62, size=11.5, color=BLUE, bold=True)
    text(slide, "MJX：金字塔近似 → 可能出现方向性伪影", 2.4, 6.08, 4.55, 0.62, size=11, color=MUTED)
    text(slide, "MJWarp：椭圆摩擦锥 → 粘滞—滑动过渡更连续", 7.0, 6.08, 5.05, 0.62, size=11, color=MUTED)
    footer(slide, 6)
    return slide


def add_branching_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "02", "ROOT CAUSE  / 根本原因", "动态分支与静态图，决定了接触处理方式")

    # Left flow: MJWarp.
    rect(slide, 0.9, 1.85, 5.52, 4.82, PALE_GREEN, None)
    text(slide, "MJWarp · 每个线程可独立决策", 1.22, 2.05, 4.9, 0.42, size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    nodes_left = [
        ("检测接触", 1.34, 2.88, 1.35, 0.68, GREEN, WHITE),
        ("运行约束求解", 5.02, 2.82, 1.1, 0.82, GREEN, WHITE),
        ("跳过", 3.41, 4.62, 1.34, 0.68, WHITE, MUTED),
    ]
    for label, x, y, w, h, fillc, tc in nodes_left:
        rect(slide, x, y, w, h, fillc, GREEN if fillc == WHITE else None)
        text(slide, label, x + 0.06, y, w - 0.12, h, size=11.5, color=tc, bold=True, align=PP_ALIGN.CENTER)
    decision = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(3.25), Inches(2.67), Inches(1.5), Inches(1.1)
    )
    set_fill(decision, WHITE)
    set_line(decision, GREEN, 1.2)
    text(slide, "接触\n活跃？", 3.52, 2.82, 0.96, 0.76, size=10.8, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    stage_arrow(slide, 2.78, 3.04, 0.32, 0.34, GREEN)
    stage_arrow(slide, 4.76, 3.04, 0.22, 0.34, GREEN)
    text(slide, "是", 4.75, 2.75, 0.32, 0.22, size=9, color=GREEN, bold=True)
    down = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(3.78), Inches(3.79), Inches(0.44), Inches(0.64)
    )
    set_fill(down, GREEN)
    set_line(down, None)
    text(slide, "否", 4.15, 3.94, 0.32, 0.22, size=9, color=MUTED, bold=True)
    rect(slide, 1.42, 5.75, 4.65, 0.58, NAVY, None)
    text(slide, "LCP / 迭代求解 → 严格非穿透", 1.42, 5.75, 4.65, 0.58, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Right flow: MJX.
    rect(slide, 6.86, 1.85, 5.54, 4.82, PALE_PURPLE, None)
    text(slide, "MJX · 所有环境共享静态计算图", 7.18, 2.05, 4.9, 0.42, size=17, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    fixed_nodes = [
        ("固定接触槽位", 7.3, 2.82, 1.55, 0.7),
        ("填充到最坏情况", 9.0, 2.82, 1.7, 0.7),
        ("连续罚函数", 10.85, 2.82, 1.15, 0.7),
    ]
    for i, (label, x, y, w, h) in enumerate(fixed_nodes):
        fillc = PURPLE if i in (0, 2) else WHITE
        rect(slide, x, y, w, h, fillc, PURPLE if fillc == WHITE else None)
        text(slide, label, x + 0.05, y, w - 0.1, h, size=11, color=WHITE if fillc != WHITE else PURPLE, bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            stage_arrow(slide, x + w + 0.08, y + 0.18, 0.25, 0.34, PURPLE)
    text(slide, "所有批次执行同一序列", 7.3, 3.83, 4.7, 0.38, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    bullet_lines(
        slide,
        ["接触为 0 也保留槽位", "连续梯度更适合优化", "复杂接触会引入填充开销"],
        7.65,
        4.32,
        4.15,
        0.48,
        dot_color=PURPLE,
        size=10.8,
    )
    rect(slide, 7.3, 5.75, 4.7, 0.58, PURPLE, None)
    text(slide, "平滑力—位移关系 → 软接触", 7.3, 5.75, 4.7, 0.58, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 7)
    return slide


def add_humanoid_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "03", "HUMANOID  / 人形机器人", "用 MJX 训练，用 MJWarp 揭示真实世界会放大的问题")

    stages = [
        (0.9, "阶段 1", "MJX 批量训练", "数千并行环境\n领域随机化\n快速策略迭代", PURPLE, PALE_PURPLE),
        (4.63, "阶段 2", "MJWarp 高保真验证", "硬接触冲击\n摩擦与滑移\n不平坦地形", GREEN, PALE_GREEN),
        (8.36, "阶段 3", "真实机器人部署", "低速安全验证\n传感器闭环\n逐步放开边界", ORANGE, PALE_ORANGE),
    ]
    for i, (x, stage, title_value, desc, color, pale) in enumerate(stages):
        rect(slide, x, 2.0, 3.18, 3.55, WHITE, color, radius=True, line_width=1.4)
        rect(slide, x + 0.25, 2.26, 0.86, 0.34, pale, None)
        text(slide, stage, x + 0.25, 2.26, 0.86, 0.34, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        circle(slide, x + 1.25, 2.78, 0.68, color)
        text(slide, str(i + 1), x + 1.25, 2.78, 0.68, 0.68, size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.3, 3.67, 2.58, 0.52, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, x + 0.35, 4.3, 2.48, 0.88, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            stage_arrow(slide, x + 3.33, 3.35, 0.78, 0.58, BLUE)

    rect(slide, 0.9, 5.9, 11.5, 0.72, NAVY, None)
    text(slide, "放行门槛", 1.18, 5.9, 1.2, 0.72, size=12, color=CYAN, bold=True)
    metrics = [
        ("落脚冲击峰值", 2.58),
        ("足端滑移距离", 5.1),
        ("关节振荡 / 抖动", 7.55),
        ("不同地面成功率", 10.0),
    ]
    for label, x in metrics:
        circle(slide, x, 6.18, 0.14, CYAN)
        text(slide, label, x + 0.25, 5.9, 2.0, 0.72, size=10.8, color=WHITE)
    footer(slide, 8)
    return slide


def add_physx_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "04", "COMPARISON  / MJWarp vs PhysX", "机器人原生的约束表达，对比通用引擎的工程成熟度")

    rect(slide, 0.9, 1.84, 3.0, 1.17, PALE_GREEN, None)
    text(slide, "MJWarp", 1.15, 1.98, 1.25, 0.38, size=18, color=GREEN, bold=True)
    text(slide, "广义坐标 · 运动学树", 1.15, 2.39, 2.45, 0.34, size=11.5, color=INK, bold=True)
    text(slide, "关节按定义不可分离", 1.15, 2.7, 2.45, 0.26, size=10.2, color=MUTED)

    rect(slide, 9.4, 1.84, 3.0, 1.17, PALE_ORANGE, None)
    text(slide, "PhysX 5", 9.65, 1.98, 1.25, 0.38, size=18, color=ORANGE, bold=True)
    text(slide, "最大坐标 · 约束连接", 9.65, 2.39, 2.45, 0.34, size=11.5, color=INK, bold=True)
    text(slide, "TGS 降低高应力下的漂移", 9.65, 2.7, 2.45, 0.26, size=10.2, color=MUTED)

    headers = ["维度", "MJWarp", "PhysX 5 / Isaac"]
    widths = [2.25, 4.0, 4.25]
    x0, y0 = 1.42, 3.34
    xx = x0
    for w, hdr, color in zip(widths, headers, [BLUE, GREEN, ORANGE]):
        rect(slide, xx, y0, w, 0.5, color, None, radius=False)
        text(slide, hdr, xx, y0, w, 0.5, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        xx += w
    rows = [
        ("关节约束", "运动学树严格保持", "TGS 约束近似保持"),
        ("接触 / 摩擦", "硬接触 · 椭圆锥", "稳定接触 · 棱锥近似"),
        ("时间步长", "较大步长仍稳定", "常需更小子步"),
        ("复杂场景生态", "机器人学习优先", "传感器 / USD / Omniverse 完整"),
    ]
    for i, row in enumerate(rows):
        yy = y0 + 0.59 + i * 0.63
        xx = x0
        for j, (w, val) in enumerate(zip(widths, row)):
            fillc = PALE if j == 0 else WHITE
            rect(slide, xx, yy, w, 0.55, fillc, BORDER, radius=False, line_width=0.7)
            text(
                slide,
                val,
                xx + 0.08,
                yy,
                w - 0.16,
                0.55,
                size=10.8,
                color=INK if j == 0 else MUTED,
                bold=(j == 0),
                align=PP_ALIGN.CENTER,
            )
            xx += w
    rect(slide, 1.42, 6.52, 10.5, 0.18, BLUE, None, radius=False)
    footer(slide, 9)
    return slide


def add_selection_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "04", "SELECTION  / 引擎选型", "先看任务形态，再看硬件与生态约束")

    x0, y0 = 0.9, 1.82
    widths = [4.35, 2.35, 2.35, 2.45]
    headers = ["任务 / 约束", "MJX", "MJWarp", "PhysX"]
    colors = [BLUE, PURPLE, GREEN, ORANGE]
    xx = x0
    for w, hdr, color in zip(widths, headers, colors):
        rect(slide, xx, y0, w, 0.55, color, None, radius=False)
        text(slide, hdr, xx, y0, w, 0.55, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        xx += w

    rows = [
        ("大规模人形 / 四足运动训练", 3, 2, 1),
        ("灵巧手、抓取、复杂网格", 1, 3, 2),
        ("精密装配 / 不平坦地形验证", 1, 3, 2),
        ("Omniverse 大场景与传感器", 1, 2, 3),
        ("端到端可微 / JAX 优化", 3, 1, 1),
        ("TPU / Mac / 跨平台", 3, 1, 1),
    ]
    for i, row in enumerate(rows):
        yy = y0 + 0.65 + i * 0.69
        fillc = PALE if i % 2 == 0 else WHITE
        xx = x0
        rect(slide, xx, yy, widths[0], 0.6, fillc, BORDER, radius=False, line_width=0.6)
        text(slide, row[0], xx + 0.22, yy, widths[0] - 0.35, 0.6, size=11.3, color=INK, bold=True)
        xx += widths[0]
        for j, score in enumerate(row[1:]):
            w = widths[j + 1]
            rect(slide, xx, yy, w, 0.6, fillc, BORDER, radius=False, line_width=0.6)
            color = [PURPLE, GREEN, ORANGE][j]
            dot_d = 0.16
            total_w = score * dot_d + (score - 1) * 0.09
            start = xx + (w - total_w) / 2
            for k in range(score):
                circle(slide, start + k * (dot_d + 0.09), yy + 0.22, dot_d, color)
            xx += w
    text(slide, "● 首选倾向越强", 9.88, 6.65, 2.5, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    footer(slide, 10)
    return slide


def add_platform_application_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(
        slide,
        "A1",
        "APPENDIX  / NVIDIA 仿真平台",
        "从 OpenUSD 场景描述到零迁移训练，形成一体化闭环",
    )

    stages = [
        ("01", "OpenUSD", "统一场景描述\n资产 · 材质 · 机器人", BLUE, PALE_BLUE),
        ("02", "Isaac Sim", "可视化与传感器\n数字场景搭建", CYAN, PALE_BLUE),
        ("03", "双物理后端", "PhysX 通用仿真\n→ MJWarp 高保真", ORANGE, PALE_ORANGE),
        ("04", "Isaac Lab / MjLab", "GPU 高并发\n强化学习训练", GREEN, PALE_GREEN),
        ("05", "零迁移训练", "复用场景与策略链路\n面向真实机器人", PURPLE, PALE_PURPLE),
    ]
    for i, (num, title_value, desc, color, pale) in enumerate(stages):
        x = 0.9 + i * 2.36
        rect(slide, x, 1.9, 2.02, 3.28, WHITE, color, radius=True, line_width=1.25)
        rect(slide, x + 0.18, 2.12, 0.55, 0.31, pale, None)
        text(
            slide,
            num,
            x + 0.18,
            2.12,
            0.55,
            0.31,
            size=9.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        circle(slide, x + 0.69, 2.72, 0.64, color)
        text(
            slide,
            title_value[0],
            x + 0.69,
            2.72,
            0.64,
            0.64,
            size=17,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(
            slide,
            title_value,
            x + 0.16,
            3.55,
            1.7,
            0.5,
            size=14.5 if i != 3 else 12.8,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(
            slide,
            desc,
            x + 0.16,
            4.15,
            1.7,
            0.68,
            size=9.8,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        if i < len(stages) - 1:
            stage_arrow(slide, x + 2.07, 3.22, 0.36, 0.46, BLUE)

    summaries = [
        ("统一场景", "OpenUSD 作为资产与数据骨架", BLUE, PALE_BLUE),
        ("物理升级", "保留 PhysX 工程能力，引入 MJWarp 接触保真度", GREEN, PALE_GREEN),
        ("训练闭环", "Isaac Lab / MjLab 承接批量训练与策略复用", ORANGE, PALE_ORANGE),
    ]
    for i, (title_value, desc, color, pale) in enumerate(summaries):
        x = 0.9 + i * 3.92
        rect(slide, x, 5.52, 3.66, 1.12, pale, None)
        circle(slide, x + 0.24, 5.78, 0.55, color)
        text(
            slide,
            str(i + 1),
            x + 0.24,
            5.78,
            0.55,
            0.55,
            size=12,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(slide, title_value, x + 0.98, 5.6, 2.4, 0.34, size=13, color=INK, bold=True)
        text(slide, desc, x + 0.98, 5.96, 2.42, 0.47, size=9.3, color=MUTED)
    footer(slide, 13)
    return slide


def add_welding_application_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(
        slide,
        "A2",
        "APPENDIX  / 焊接与传动系统",
        "焊接机器人、减速机与伺服电机：三类 GPU 仿真能力",
    )

    cards = [
        (
            0.9,
            "NV Warp 多体动力学",
            "机器人—减速机—伺服电机\n传动链的高并发仿真",
            "动力学 / 接触 / 摩擦 / 载荷",
            GREEN,
            PALE_GREEN,
            "D",
        ),
        (
            4.82,
            "GPU 加速 3D 视觉",
            "在 Isaac Sim 中批量生成\n相机、深度与点云观测",
            "工件定位 / 焊缝感知 / 视觉闭环",
            CYAN,
            PALE_BLUE,
            "V",
        ),
        (
            8.74,
            "多体 + 热力学耦合",
            "将机器人运动与焊接热场\n放入统一 GPU 仿真流程",
            "轨迹 / 动力学 / 热输入 / 温度场",
            ORANGE,
            PALE_ORANGE,
            "T",
        ),
    ]
    for x, title_value, desc, tags, color, pale, initial in cards:
        rect(slide, x, 1.88, 3.66, 2.32, WHITE, color, radius=True, line_width=1.3)
        circle(slide, x + 0.25, 2.17, 0.64, color)
        text(
            slide,
            initial,
            x + 0.25,
            2.17,
            0.64,
            0.64,
            size=17,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(slide, title_value, x + 1.05, 2.05, 2.32, 0.45, size=16, color=INK, bold=True)
        text(slide, desc, x + 0.3, 2.92, 3.06, 0.62, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)
        rect(slide, x + 0.32, 3.67, 3.02, 0.33, pale, None)
        text(
            slide,
            tags,
            x + 0.38,
            3.67,
            2.9,
            0.33,
            size=8.8,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    rect(slide, 0.9, 4.58, 11.5, 1.72, PALE, None)
    text(slide, "统一数字样机闭环", 1.16, 4.76, 2.1, 0.38, size=14, color=BLUE, bold=True)
    flow = [
        ("OpenUSD", "工件 / 机器人 / 场景", BLUE),
        ("Isaac Sim", "可视化 / 3D 视觉", CYAN),
        ("NV Warp", "机器人与传动链", GREEN),
        ("GPU 热力学", "焊接热场", ORANGE),
        ("MjLab", "控制 / RL / 验证", PURPLE),
    ]
    for i, (title_value, desc, color) in enumerate(flow):
        x = 1.16 + i * 2.19
        rect(slide, x, 5.27, 1.78, 0.77, WHITE, color, radius=True, line_width=1)
        text(slide, title_value, x + 0.06, 5.31, 1.66, 0.28, size=10.8, color=color, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, x + 0.06, 5.59, 1.66, 0.28, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            stage_arrow(slide, x + 1.87, 5.47, 0.27, 0.34, BLUE)
    rect(slide, 0.9, 6.5, 11.5, 0.19, BLUE, None, radius=False)
    footer(slide, 14)
    return slide


def add_thermal_application_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(
        slide,
        "A3",
        "APPENDIX  / THERMAL SIMULATION",
        "热仿真：从移动热源到温度场反演",
    )

    rect(slide, 0.9, 1.84, 3.45, 4.7, NAVY, None)
    text(
        slide,
        "∂T/∂t",
        1.16,
        2.22,
        2.92,
        0.72,
        size=35,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_MONO,
    )
    text(
        slide,
        "= α∇²T + Q/(ρc)",
        1.05,
        2.94,
        3.15,
        0.58,
        size=21,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_MONO,
    )
    line(slide, 1.38, 3.75, 3.87, 3.75, NAVY_LINE, 1)
    text(
        slide,
        "Warp 的角色",
        1.22,
        3.96,
        2.8,
        0.38,
        size=14,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    bullet_lines(
        slide,
        [
            "GPU 并行温度更新",
            "网格 / 网格单元积分",
            "批量参数扫描",
            "自动微分与反演",
        ],
        1.4,
        4.45,
        2.55,
        0.44,
        dot_color=CYAN,
        size=10.4,
        text_color=LIGHT_TEXT,
    )

    steps = [
        ("01", "离散域", "Grid / Mesh / NanoVDB\n材料参数 k · ρ · c", BLUE, PALE_BLUE),
        ("02", "热源与边界", "移动热源 Q(x,t)\n对流 · 辐射 · 定温", ORANGE, PALE_ORANGE),
        ("03", "GPU 求解", "Warp kernel / warp.fem\n显式步进或线性求解", GREEN, PALE_GREEN),
        ("04", "批量与反演", "工艺参数扫描\n梯度校准 / 优化", PURPLE, PALE_PURPLE),
    ]
    for i, (num, title_value, desc, color, pale) in enumerate(steps):
        col, row = i % 2, i // 2
        x = 4.7 + col * 3.86
        y = 1.84 + row * 1.72
        rect(slide, x, y, 3.54, 1.48, WHITE, BORDER, radius=True)
        rect(slide, x + 0.2, y + 0.24, 0.58, 0.58, color, None, radius=False)
        text(slide, num, x + 0.2, y + 0.24, 0.58, 0.58, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.98, y + 0.15, 2.25, 0.4, size=14.5, color=INK, bold=True)
        text(slide, desc, x + 0.98, y + 0.58, 2.28, 0.62, size=9.6, color=MUTED)
        rect(slide, x + 0.98, y + 1.19, 2.25, 0.12, pale, None, radius=False)

    rect(slide, 4.7, 5.42, 7.55, 1.12, PALE, None)
    use_cases = [
        ("焊接热场", "Goldak / 移动热源"),
        ("伺服与减速机", "损耗—温升—热漂移"),
        ("逆问题", "热参数辨识与工艺优化"),
    ]
    for i, (title_value, desc) in enumerate(use_cases):
        x = 4.98 + i * 2.4
        circle(slide, x, 5.72, 0.47, [ORANGE, GREEN, PURPLE][i])
        text(slide, str(i + 1), x, 5.72, 0.47, 0.47, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.62, 5.52, 1.55, 0.31, size=10.8, color=INK, bold=True)
        text(slide, desc, x + 0.62, 5.86, 1.58, 0.31, size=8.7, color=MUTED)
    text(
        slide,
        "能力边界：Warp 提供 GPU kernel、FEM 与自动微分；材料模型、热源模型和耦合策略仍需按场景实现。",
        0.9,
        6.78,
        11.5,
        0.18,
        size=7.2,
        color="8B95A5",
    )
    footer(slide, 15)
    return slide


def add_electromagnetic_application_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(
        slide,
        "A4",
        "APPENDIX  / ELECTROMAGNETICS",
        "电磁仿真：从 Maxwell 离散到磁场与力矩优化",
    )

    rect(slide, 0.9, 1.84, 5.52, 3.72, PALE_BLUE, None)
    text(slide, "路径 1 · FDTD / 自定义 Kernel", 1.22, 2.05, 4.88, 0.42, size=16.5, color=BLUE, bold=True)
    rich_text(
        slide,
        [
            {"text": "∂E/∂t", "bold": True, "color": BLUE, "font": FONT_MONO, "size": 19},
            {"text": "  ↔  ", "bold": True, "color": MUTED, "size": 17},
            {"text": "∇×H", "bold": True, "color": BLUE, "font": FONT_MONO, "size": 19},
        ],
        1.22,
        2.64,
        4.85,
        0.52,
        align=PP_ALIGN.CENTER,
    )
    flow_left = [
        ("交错网格", "E / H 场"),
        ("时间步进", "curl 更新"),
        ("边界处理", "PML / 激励"),
    ]
    for i, (title_value, desc) in enumerate(flow_left):
        x = 1.2 + i * 1.62
        rect(slide, x, 3.43, 1.34, 0.86, WHITE, BLUE, radius=True, line_width=0.9)
        text(slide, title_value, x + 0.05, 3.46, 1.24, 0.32, size=10.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, x + 0.05, 3.82, 1.24, 0.29, size=8.8, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            stage_arrow(slide, x + 1.4, 3.68, 0.2, 0.3, BLUE)
    text(
        slide,
        "适合：瞬态电磁波、规则网格、大批量设计点",
        1.22,
        4.72,
        4.9,
        0.42,
        size=10.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    rect(slide, 6.86, 1.84, 5.54, 3.72, PALE_PURPLE, None)
    text(slide, "路径 2 · warp.fem / Curl-Curl", 7.18, 2.05, 4.88, 0.42, size=16.5, color=PURPLE, bold=True)
    rich_text(
        slide,
        [
            {"text": "∇×(μ⁻¹∇×A)", "bold": True, "color": PURPLE, "font": FONT_MONO, "size": 18},
            {"text": " = J", "bold": True, "color": MUTED, "font": FONT_MONO, "size": 18},
        ],
        7.18,
        2.64,
        4.86,
        0.52,
        align=PP_ALIGN.CENTER,
    )
    flow_right = [
        ("非结构网格", "复杂几何"),
        ("Nédélec 空间", "curl 相容"),
        ("稀疏求解", "磁场 / 力"),
    ]
    for i, (title_value, desc) in enumerate(flow_right):
        x = 7.16 + i * 1.62
        rect(slide, x, 3.43, 1.34, 0.86, WHITE, PURPLE, radius=True, line_width=0.9)
        text(slide, title_value, x + 0.05, 3.46, 1.24, 0.32, size=9.9, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, x + 0.05, 3.82, 1.24, 0.29, size=8.8, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            stage_arrow(slide, x + 1.4, 3.68, 0.2, 0.3, PURPLE)
    text(
        slide,
        "官方示例：2D magnetostatics · curl-curl formulation",
        7.18,
        4.72,
        4.86,
        0.42,
        size=10.2,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    use_cases = [
        ("伺服电机", "磁场 / 力矩 / 参数优化", GREEN),
        ("感应加热与焊接", "电磁—热源耦合", ORANGE),
        ("传感器与执行器", "场分布 / 灵敏度设计", CYAN),
    ]
    for i, (title_value, desc, color) in enumerate(use_cases):
        x = 0.9 + i * 3.92
        rect(slide, x, 5.82, 3.66, 0.76, WHITE, BORDER, radius=True)
        circle(slide, x + 0.2, 5.98, 0.44, color)
        text(slide, str(i + 1), x + 0.2, 5.98, 0.44, 0.44, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.78, 5.82, 1.35, 0.34, size=10.7, color=INK, bold=True)
        text(slide, desc, x + 2.03, 5.82, 1.4, 0.34, size=8.7, color=MUTED, align=PP_ALIGN.RIGHT)
    text(
        slide,
        "能力边界：Warp 官方提供磁静力 FEM 示例；全波、频域、多材料与 PML 等能力需自行构建和验证。",
        0.9,
        6.78,
        11.5,
        0.18,
        size=7.2,
        color="8B95A5",
    )
    footer(slide, 16)
    return slide


def add_fea_application_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(
        slide,
        "A5",
        "APPENDIX  / FINITE ELEMENT ANALYSIS",
        "有限元分析：GPU 加速装配、求解与可微设计",
    )

    stages = [
        ("01", "Geometry", "Grid / Mesh\nNanoVDB", BLUE),
        ("02", "Function Space", "P / Q / S\n标量与向量场", CYAN),
        ("03", "Integrate", "线性 / 双线性形式\n数值积分", ORANGE),
        ("04", "Sparse System", "BSR 矩阵\n边界条件", GREEN),
        ("05", "Solve + Grad", "CG / 迭代求解\n自动微分", PURPLE),
    ]
    for i, (num, title_value, desc, color) in enumerate(stages):
        x = 0.9 + i * 2.36
        rect(slide, x, 1.86, 2.02, 2.36, WHITE, color, radius=True, line_width=1.2)
        rect(slide, x + 0.18, 2.09, 0.55, 0.31, PALE, None)
        text(slide, num, x + 0.18, 2.09, 0.55, 0.31, size=9.3, color=color, bold=True, align=PP_ALIGN.CENTER)
        circle(slide, x + 0.7, 2.58, 0.62, color)
        text(slide, title_value[0], x + 0.7, 2.58, 0.62, 0.62, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.15, 3.38, 1.72, 0.4, size=12.5 if i != 1 else 10.8, color=INK, bold=True, align=PP_ALIGN.CENTER)
        text(slide, desc, x + 0.16, 3.78, 1.7, 0.34, size=8.8, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            stage_arrow(slide, x + 2.07, 2.86, 0.36, 0.44, BLUE)

    rect(slide, 0.9, 4.58, 7.52, 1.96, PALE, None)
    text(slide, "warp.fem 可直接覆盖的典型问题", 1.18, 4.77, 6.9, 0.36, size=14, color=BLUE, bold=True)
    fem_cases = [
        ("扩散 / 对流", "热传导与输运", BLUE),
        ("弹性力学", "应力、应变、变形", GREEN),
        ("接触", "非匹配网格接触", ORANGE),
        ("优化", "弹性体形状优化", PURPLE),
    ]
    for i, (title_value, desc, color) in enumerate(fem_cases):
        x = 1.18 + (i % 2) * 3.42
        y = 5.3 + (i // 2) * 0.52
        circle(slide, x, y + 0.1, 0.18, color)
        text(slide, title_value, x + 0.32, y, 1.15, 0.35, size=10.3, color=INK, bold=True)
        text(slide, desc, x + 1.52, y, 1.55, 0.35, size=9, color=MUTED)

    rect(slide, 8.72, 4.58, 3.68, 1.96, NAVY, None)
    text(slide, "机器人结构应用", 9.0, 4.79, 3.1, 0.36, size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    bullet_lines(
        slide,
        [
            "连杆与末端刚度",
            "减速机壳体应力",
            "热—结构耦合变形",
            "轻量化 / 逆向设计",
        ],
        9.08,
        5.22,
        2.82,
        0.29,
        dot_color=CYAN,
        size=9.3,
        text_color=WHITE,
    )
    text(
        slide,
        "依据 NVIDIA Warp 文档：warp.fem 支持 PDE 的 Galerkin 离散、网格/函数空间、形式积分、稀疏系统与迭代求解。",
        0.9,
        6.78,
        11.5,
        0.18,
        size=7.2,
        color="8B95A5",
    )
    footer(slide, 17)
    return slide


def add_mjlab_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    section_header(slide, "05", "MJLAB  / 能力路线", "平台落地需要补齐八项工程能力")

    items = [
        ("01", "快速 CUDA 仿真", "构建更强的 Isaac Lab 训练底座", BLUE, PALE_BLUE),
        ("02", "执行器模型", "电机、减速器与控制延迟更贴近实机", PURPLE, PALE_PURPLE),
        ("03", "脚部碰撞模型", "聚焦落脚冲击、摩擦与滑移", GREEN, PALE_GREEN),
        ("04", "Warp → Newton", "为 Isaac Lab 的新物理后端预留接口", ORANGE, PALE_ORANGE),
        ("05", "物理保真度", "建立 MJX / MJWarp / 实机一致性指标", GREEN, PALE_GREEN),
        ("06", "训练检查点", "可恢复、可回放、可跨后端复现", BLUE, PALE_BLUE),
        ("07", "MuJoCo 内训练", "缩短模型、策略与验证之间的距离", PURPLE, PALE_PURPLE),
        ("08", "IRL 与重定向", "人形奖励学习 → 双足动作迁移", ORANGE, PALE_ORANGE),
    ]
    for i, (num, title_value, desc, color, pale) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.9 + col * 5.88
        y = 1.78 + row * 1.18
        rect(slide, x, y, 5.62, 1.0, WHITE, BORDER, radius=True)
        rect(slide, x + 0.18, y + 0.18, 0.64, 0.64, color, None, radius=False)
        text(slide, num, x + 0.18, y + 0.18, 0.64, 0.64, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 1.02, y + 0.1, 4.25, 0.4, size=14.5, color=INK, bold=True)
        text(slide, desc, x + 1.02, y + 0.5, 4.25, 0.33, size=10.3, color=MUTED)
    rect(slide, 0.9, 6.53, 11.5, 0.2, BLUE, None, radius=False)
    footer(slide, 11)
    return slide


def add_recommendation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Match the template's airy closing composition.
    circle(slide, 9.95, 3.7, 4.5, PALE_BLUE)
    circle(slide, -1.6, -1.8, 4.2, PALE)
    rect(slide, 0.9, 0.88, 0.55, 0.09, BLUE, None, radius=False)
    text(slide, "RECOMMENDATION  / 建议", 0.9, 1.08, 6.5, 0.32, size=11, color=BLUE, bold=True)
    text(slide, "平台、物理、应用分层", 0.88, 1.55, 8.2, 0.72, size=34, color=INK, bold=True)
    text(slide, "用双试点验证 NVIDIA 仿真闭环", 0.9, 2.28, 8.4, 0.5, size=21, color=BLUE, bold=True)

    recommendations = [
        ("01", "OpenUSD 统一场景", "让资产、传感器与仿真数据跨阶段复用", BLUE),
        ("02", "双物理后端分工", "PhysX 承接通用工程，MJWarp 承接高保真与并发", GREEN),
        ("03", "双应用试点", "新型机器人零迁移训练 + 焊接机器人多物理场", ORANGE),
    ]
    for i, (num, title_value, desc, color) in enumerate(recommendations):
        yy = 3.12 + i * 1.03
        circle(slide, 1.0, yy + 0.15, 0.52, color)
        text(slide, num, 1.0, yy + 0.15, 0.52, 0.52, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title_value, 1.75, yy, 3.0, 0.4, size=15, color=INK, bold=True)
        text(slide, desc, 4.65, yy, 5.75, 0.4, size=11, color=MUTED)
        if i < 2:
            line(slide, 1.75, yy + 0.76, 9.8, yy + 0.76, BORDER, 0.8)

    rect(slide, 0.95, 6.52, 0.6, 0.08, BLUE, None, radius=False)
    text(
        slide,
        "下一步：以同一 OpenUSD 场景建立 Isaac Sim → MJWarp / MjLab → 实机回归基线",
        0.95,
        6.67,
        10.6,
        0.34,
        size=12,
        color=MUTED,
    )
    text(slide, "12", 11.9, 7.08, 0.5, 0.22, size=8, color="9AA4B5", align=PP_ALIGN.RIGHT)
    return slide


def build():
    prs = Presentation(TEMPLATE)
    wipe_template_slides(prs)
    prs.core_properties.title = "MJWarp · MJX · PhysX：机器人仿真后端对比、平台应用与 MjLab 选型建议"
    prs.core_properties.subject = "基于 MJ Warp, MJX, Physx.md，使用指定 NVIDIA 模板制作"
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.keywords = "MJWarp, MJX, PhysX, MuJoCo, robotics, simulation, MjLab"
    prs.core_properties.comments = "All claims and recommendations are synthesized from the supplied Markdown source."

    add_title_slide(prs)
    add_contents_slide(prs)
    add_executive_slide(prs)
    add_architecture_slide(prs)
    add_mjx_mjwarp_slide(prs)
    add_contact_slide(prs)
    add_branching_slide(prs)
    add_humanoid_slide(prs)
    add_physx_slide(prs)
    add_selection_slide(prs)
    add_mjlab_slide(prs)
    add_recommendation_slide(prs)
    add_platform_application_slide(prs)
    add_welding_application_slide(prs)
    add_thermal_application_slide(prs)
    add_electromagnetic_application_slide(prs)
    add_fea_application_slide(prs)

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
